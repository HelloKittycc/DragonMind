from io import BytesIO

from src.modules.knowledge.parsers.base import KnowledgeFileParser, ParserError


class PptxKnowledgeParser(KnowledgeFileParser):
    supported_extensions = {".pptx"}

    def extract_text(self, content: bytes, filename: str) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserError("PPTX parser dependency is not installed") from exc

        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            raise ParserError("PPTX file could not be parsed") from exc

        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            if texts:
                slides.append(f"[Slide {index}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
